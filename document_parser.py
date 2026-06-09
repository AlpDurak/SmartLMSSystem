import fitz  # PyMuPDF
from pptx import Presentation
import os

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    return "".join(text_parts)

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
    return "\n".join(text_parts) + "\n" if text_parts else ""

def get_document_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        return ""
